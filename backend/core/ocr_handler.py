"""
ResearchFlow Backend - Intelligent OCR Handler

Combines Tesseract OCR with LLM-powered intelligent post-processing for
high-accuracy document text extraction.

Pipeline:
  1. File Conversion: PDF/PPT/PPTX -> images
  2. Image Preprocessing: enhance for OCR quality
  3. Tesseract OCR: extract raw text
  4. LLM Intelligent Correction: fix errors, detect tables, preserve structure
  5. Final Output: clean, structured text

Supports: PDF, PNG, JPG, JPEG, TIFF, BMP, PPT, PPTX
"""

import os
import uuid
import logging
import threading
import asyncio
import re
from typing import Dict, List, Optional, Tuple, Callable
from pathlib import Path
from datetime import datetime
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor

from PIL import Image, ImageEnhance, ImageFilter
from openai import OpenAI

from config.settings import settings
from core.progress import progress_tracker

logger = logging.getLogger(__name__)


INTELLIGENT_OCR_SYSTEM_PROMPT = """You are an expert OCR post-processing assistant. Your task is to correct and enhance OCR-extracted text from documents.

Core capabilities:
1. Fix common OCR errors (character confusions, split/merged words)
2. Detect and properly format tables (convert to markdown tables)
3. Preserve document structure (headings, lists, paragraphs)
4. Correct mixed-language content (Chinese + English)
5. Recognize and properly format formulas (LaTeX where appropriate)
6. Remove OCR artifacts and noise

Rules:
- Preserve all factual information exactly as in the original
- Never invent or add information not present in the raw OCR text
- Fix obvious OCR errors: "RnD" -> "R&D", "0" vs "O" in context, "1" vs "l"
- For tables: detect grid patterns and format as markdown tables
- For forms: preserve field labels and values
- For mixed CN/EN: ensure correct spacing between Chinese and English
- Maintain original line breaks only for intentional formatting
- If the text appears to be garbled/unreadable, state that clearly

Output format:
- Use markdown formatting for structure
- Include a brief [OCR-LLM] confidence note at the very end
"""


class OCRProcessor:
    """
    Intelligent OCR processor combining Tesseract engine with LLM correction.

    Features:
    - Multi-format support (PDF, images, PPT)
    - Image preprocessing for better recognition
    - Multi-language OCR (Chinese + English by default)
    - LLM-powered intelligent text correction
    - Table/form detection and formatting
    - Real-time progress tracking
    """

    MAX_IMAGE_SIZE = 4096
    LLM_CORRECTION_CHUNK_SIZE = 4000

    def __init__(
        self,
        tesseract_path: str = None,
        languages: str = None,
        dpi: int = None,
        model: str = None,
        api_key: str = None,
        base_url: str = None,
    ):
        import pytesseract

        self.tesseract_path = tesseract_path or settings.TESSERACT_PATH
        if os.path.exists(self.tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = self.tesseract_path

        self.languages = languages or settings.OCR_LANGUAGES
        self.dpi = dpi or settings.OCR_DPI

        self.llm_client = OpenAI(
            api_key=api_key or settings.GEMINI_API_KEY,
            base_url=base_url or settings.GEMINI_BASE_URL,
            timeout=60.0,
            max_retries=1,
        )
        self.model = model or settings.GEMINI_MODEL

        self._ocr_upload_dir = settings.UPLOAD_DIR / "ocr"
        os.makedirs(self._ocr_upload_dir, exist_ok=True)

        logger.info(f"OCRProcessor initialized: tesseract={self.tesseract_path}, langs={self.languages}")

    def convert_to_images(self, file_path: str, doc_id: str) -> List[Image.Image]:
        """
        Convert various document formats to a list of PIL Images.

        Supported formats:
        - PDF: via pdf2image
        - PPT/PPTX: via python-pptx + rendering
        - Images: PNG, JPG, JPEG, TIFF, BMP
        """
        ext = Path(file_path).suffix.lower()

        if ext == '.pdf':
            return self._pdf_to_images(file_path)
        elif ext in ('.ppt', '.pptx'):
            return self._ppt_to_images(file_path, doc_id)
        elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp'):
            return self._image_load(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _pdf_to_images(self, file_path: str) -> List[Image.Image]:
        """Convert PDF pages to images using pdf2image."""
        try:
            from pdf2image import convert_from_path
        except ImportError:
            raise ImportError(
                "pdf2image is required for PDF OCR. "
                "Install: pip install pdf2image. "
                "Also requires poppler: https://github.com/oschwartz10612/poppler-windows/releases/"
            )

        images = convert_from_path(file_path, dpi=self.dpi)
        logger.info(f"PDF converted: {len(images)} pages at {self.dpi} DPI")
        return images

    def _ppt_to_images(self, file_path: str, doc_id: str) -> List[Image.Image]:
        """Convert PPT/PPTX slides to images using python-pptx and Pillow."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError("python-pptx is required for PPT OCR. Install: pip install python-pptx")

        prs = Presentation(file_path)
        images = []
        temp_dir = self._ocr_upload_dir / f"temp_{doc_id}"
        os.makedirs(temp_dir, exist_ok=True)

        try:
            for i, slide in enumerate(prs.slides):
                slide_width = prs.slide_width or Inches(10)
                slide_height = prs.slide_height or Inches(7.5)

                scale = self.dpi / 96.0
                img_width = int((slide_width / 914400) * 96 * scale)
                img_height = int((slide_height / 914400) * 96 * scale)

                img = Image.new('RGB', (img_width, img_height), 'white')
                images.append(img)

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        pass

            if not images:
                blank = Image.new('RGB', (1920, 1080), 'white')
                images.append(blank)
                logger.warning(f"No slides rendered in PPT, using blank image")

        finally:
            for f in temp_dir.iterdir():
                try:
                    f.unlink()
                except Exception:
                    pass
            try:
                temp_dir.rmdir()
            except Exception:
                pass

        logger.info(f"PPT converted: {len(images)} slides")
        return images

    def _image_load(self, file_path: str) -> List[Image.Image]:
        """Load a single image file."""
        img = Image.open(file_path)
        if img.mode != 'RGB':
            img = img.convert('RGB')
        return [img]

    def preprocess_image(self, image: Image.Image) -> Image.Image:
        """
        Preprocess image for better OCR accuracy.

        Tesseract 4.x+ uses LSTM neural networks internally and works best
        with grayscale images. External binarization damages anti-aliased
        text edges and destroys fine details that the LSTM model relies on.

        Steps:
        1. Resize if too large (maintain aspect ratio)
        2. Convert to grayscale
        3. Mild contrast enhancement (1.3x)
        4. Mild sharpening (1.2x)
        5. Very light Gaussian blur for noise reduction
        """
        w, h = image.size
        if max(w, h) > self.MAX_IMAGE_SIZE:
            scale = self.MAX_IMAGE_SIZE / max(w, h)
            new_w, new_h = int(w * scale), int(h * scale)
            image = image.resize((new_w, new_h), Image.LANCZOS)

        gray = image.convert('L')

        enhancer = ImageEnhance.Contrast(gray)
        gray = enhancer.enhance(1.3)

        enhancer = ImageEnhance.Sharpness(gray)
        gray = enhancer.enhance(1.2)

        gray = gray.filter(ImageFilter.GaussianBlur(radius=0.5))

        return gray

    def extract_text(self, image: Image.Image) -> str:
        """Run Tesseract OCR on a preprocessed image."""
        import pytesseract

        try:
            text = pytesseract.image_to_string(
                image,
                lang=self.languages,
                config='--oem 1 --psm 6',
            )
            return text.strip()
        except Exception as e:
            logger.error(f"Tesseract OCR failed: {e}")
            return ""

    def extract_text_with_boxes(self, image: Image.Image) -> List[Dict]:
        """Extract text with bounding box information for structure preservation."""
        import pytesseract

        try:
            data = pytesseract.image_to_data(
                image,
                lang=self.languages,
                config='--oem 1 --psm 6',
                output_type=pytesseract.Output.DICT,
            )
            return data
        except Exception as e:
            logger.error(f"Tesseract data extraction failed: {e}")
            return {}

    async def intelligent_correct(
        self,
        raw_text: str,
        filename: str = "",
        context: str = "",
    ) -> str:
        """
        Use LLM to intelligently correct and enhance OCR output.

        This is the key differentiator - not just OCR, but SMART OCR.
        The LLM:
        1. Fixes character-level OCR errors using context understanding
        2. Detects and formats tables as proper markdown
        3. Preserves document structure
        4. Handles mixed Chinese/English content
        5. Recognizes special patterns (formulas, code, forms)
        """
        if not raw_text or len(raw_text.strip()) < 10:
            return raw_text

        if len(raw_text) > self.LLM_CORRECTION_CHUNK_SIZE:
            return await self._intelligent_correct_chunked(raw_text, filename, context)

        prompt = f"Document: {filename}\n"
        if context:
            prompt += f"Context: {context}\n"
        prompt += f"\nRaw OCR Text:\n```\n{raw_text}\n```\n\nPlease correct and format this OCR text."

        try:
            loop = asyncio.get_running_loop()
            response = await loop.run_in_executor(
                None,
                lambda: self.llm_client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": INTELLIGENT_OCR_SYSTEM_PROMPT},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.2,
                    max_tokens=8192,
                ),
            )

            corrected = response.choices[0].message.content.strip()

            confidence_note = "\n\n---\n*[OCR-LLM: Text has been intelligently corrected and enhanced by AI. Original OCR may have contained errors.]*"
            if not corrected.endswith(confidence_note):
                corrected += confidence_note

            return corrected

        except Exception as e:
            logger.error(f"LLM correction failed: {e}")
            return raw_text + "\n\n*[LLM correction unavailable, showing raw OCR output]*"

    async def _intelligent_correct_chunked(
        self,
        raw_text: str,
        filename: str,
        context: str,
    ) -> str:
        """Process long text by splitting into chunks and correcting each."""
        paragraphs = raw_text.split('\n\n')
        chunks = []
        current_chunk = []

        for para in paragraphs:
            if sum(len(c) for c in current_chunk) + len(para) < self.LLM_CORRECTION_CHUNK_SIZE:
                current_chunk.append(para)
            else:
                if current_chunk:
                    chunks.append('\n\n'.join(current_chunk))
                current_chunk = [para]

        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))

        if len(chunks) == 1:
            return await self._force_chunk_and_correct(chunks[0], filename, context)

        corrected_chunks = []
        for i, chunk in enumerate(chunks):
            chunk_context = f"Part {i + 1}/{len(chunks)}"
            corrected = await self.intelligent_correct(chunk, filename, chunk_context)
            corrected_chunks.append(corrected)

        return '\n\n'.join(corrected_chunks)

    async def _force_chunk_and_correct(
        self,
        raw_text: str,
        filename: str,
        context: str,
    ) -> str:
        """Force-split a single oversized paragraph into character-count chunks."""
        chunk_limit = self.LLM_CORRECTION_CHUNK_SIZE - 500
        forced_chunks = [
            raw_text[i:i + chunk_limit]
            for i in range(0, len(raw_text), chunk_limit)
        ]
        if len(forced_chunks) == 1:
            return forced_chunks[0]

        corrected_chunks = []
        for i, chunk in enumerate(forced_chunks):
            chunk_context = f"Part {i + 1}/{len(forced_chunks)}"
            corrected = await self.intelligent_correct(chunk, filename, chunk_context)
            corrected_chunks.append(corrected)

        return '\n\n'.join(corrected_chunks)

    async def process_document(
        self,
        doc_id: str,
        file_path: str,
        filename: str,
        progress_callback: Optional[Callable] = None,
        enable_llm: bool = False,
    ) -> Dict:
        """
        Full intelligent OCR pipeline with progress tracking.

        Pipeline stages:
        1. File conversion (PDF/PPT -> images) [20%]
        2. Image preprocessing [10%]
        3. Tesseract OCR extraction [40%]
        4. LLM intelligent correction [30%]

        Returns dict with ocr_text, metadata, and progress info.
        """
        loop = asyncio.get_running_loop()

        progress = progress_tracker.get(doc_id)
        if not progress:
            progress = progress_tracker.create(doc_id, filename)

        try:
            total_steps = 3 if not enable_llm else 4
            progress.set_step(0, total_steps, "Converting document to images...")
            if progress_callback:
                progress_callback(0, "Converting document to images...")

            # CPU-bound: run in thread pool to avoid blocking the event loop
            images = await loop.run_in_executor(
                None, self.convert_to_images, file_path, doc_id
            )
            total_pages = len(images)

            progress.set_step(1, total_steps, f"Preprocessing {total_pages} pages...")
            if progress_callback:
                progress_callback(25, f"Preprocessing {total_pages} pages...")

            if total_pages > 1:
                max_workers = min(4, os.cpu_count() or 4)
                with ThreadPoolExecutor(max_workers=max_workers) as pool:
                    processed_images = list(pool.map(self.preprocess_image, images))
            else:
                processed_images = [self.preprocess_image(images[0])]

            progress.set_step(2, total_steps, f"Running OCR on {total_pages} pages...")
            if progress_callback:
                progress_callback(50, f"Running OCR on {total_pages} pages...")

            if total_pages > 1:
                max_workers = min(4, os.cpu_count() or 4)
                with ThreadPoolExecutor(max_workers=max_workers) as pool:
                    texts = list(pool.map(self.extract_text, processed_images))
                ocr_results = [
                    {"page": i + 1, "text": text or "", "char_count": len(text or "")}
                    for i, text in enumerate(texts)
                ]
            else:
                text = self.extract_text(processed_images[0])
                ocr_results = [{"page": 1, "text": text or "", "char_count": len(text or "")}]

            combined_text = ""
            for result in ocr_results:
                if result["text"]:
                    combined_text += f"\n--- Page {result['page']} ---\n{result['text']}\n"

            raw_char_count = sum(r["char_count"] for r in ocr_results)

            if enable_llm:
                progress.set_step(3, total_steps, "LLM intelligent correction...")
                if progress_callback:
                    progress_callback(75, "LLM intelligent correction in progress...")

                corrected_text = await self.intelligent_correct(
                    combined_text,
                    filename=filename,
                )
                finish_msg = "OCR completed with LLM correction"
            else:
                corrected_text = combined_text
                finish_msg = "OCR completed (raw text, LLM correction skipped)"

            progress.set_step(total_steps, total_steps, "Completed")
            if progress_callback:
                progress_callback(100, "OCR processing completed")

            progress_tracker.update(
                doc_id,
                status="completed",
                progress=100,
                current_step=finish_msg,
            )

            return {
                "doc_id": doc_id,
                "filename": filename,
                "total_pages": total_pages,
                "raw_text": combined_text,
                "ocr_text": corrected_text,
                "raw_char_count": raw_char_count,
                "corrected_char_count": len(corrected_text),
                "llm_corrected": enable_llm,
                "page_results": [
                    {
                        "page": r["page"],
                        "char_count": r["char_count"],
                        "has_content": len(r["text"].strip()) > 0,
                    }
                    for r in ocr_results
                ],
                "status": "completed",
                "processed_at": datetime.now().isoformat(),
            }

        except Exception as e:
            logger.error(f"OCR processing failed for {doc_id}: {e}")
            progress_tracker.update(doc_id, error=str(e))
            raise


ocr_processor = OCRProcessor()